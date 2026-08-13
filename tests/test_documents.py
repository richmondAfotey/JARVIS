"""Tests for the Phase 12 document-intelligence tool."""

import pytest

from tools import build_default_registry
from tools.base import ToolError
from tools.documents import ReadDocumentTool

PLAIN = {"path": None}
_DOCX = "sample.docx"
_XLSX = "sample.xlsx"
_PPTX = "sample.pptx"
_PDF = "sample.pdf"


def _minimal_pdf(text: str) -> bytes:
    content = f"BT /F1 24 Tf 100 700 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>"
        ),
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, 1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n"
    ).encode()
    return bytes(out)


# -- plain text -----------------------------------------------------------

def test_read_plain_text(tmp_path):
    target = tmp_path / "note.txt"
    target.write_text("The quick brown fox.", encoding="utf-8")
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "quick brown fox" in out


def test_read_markdown(tmp_path):
    target = tmp_path / "doc.md"
    target.write_text("# Heading\nBody text.", encoding="utf-8")
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "# Heading" in out


def test_read_csv_as_text(tmp_path):
    target = tmp_path / "data.csv"
    target.write_text("name,age\nAda,36\nGrace,45", encoding="utf-8")
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "Ada" in out
    assert "Grace" in out


def test_read_text_truncates(tmp_path):
    target = tmp_path / "long.txt"
    target.write_text("x" * 20_000, encoding="utf-8")
    out = ReadDocumentTool().execute({"path": str(target), "max_chars": 100})
    assert "truncated" in out
    assert len(out) < 500


# -- .docx ----------------------------------------------------------------

def test_read_docx(tmp_path):
    import docx

    target = tmp_path / _DOCX
    document = docx.Document()
    document.add_paragraph("Meeting notes")
    document.add_paragraph("Decision: ship it.")
    document.save(str(target))
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "Meeting notes" in out
    assert "ship it" in out


# -- .xlsx ----------------------------------------------------------------

def test_read_xlsx(tmp_path):
    from openpyxl import Workbook

    target = tmp_path / _XLSX
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Month", "Total"])
    sheet.append(["Jan", 120])
    sheet.append(["Feb", 85])
    workbook.save(str(target))
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "Sales" in out
    assert "Jan" in out
    assert "120" in out


# -- .pptx ----------------------------------------------------------------

def test_read_pptx(tmp_path):
    from pptx import Presentation

    target = tmp_path / _PPTX
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Phase 12 today"
    presentation.save(str(target))
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "Roadmap" in out
    assert "Phase 12" in out


# -- .pdf -----------------------------------------------------------------

def test_read_pdf(tmp_path):
    target = tmp_path / _PDF
    target.write_bytes(_minimal_pdf("Hello World PDF"))
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "Hello World PDF" in out


def test_read_pdf_no_text_layer(tmp_path):
    target = tmp_path / _PDF
    target.write_bytes(_minimal_pdf(""))
    out = ReadDocumentTool().execute({"path": str(target)})
    assert "no extractable text" in out


# -- errors ---------------------------------------------------------------

def test_read_missing_file(tmp_path):
    try:
        ReadDocumentTool().execute({"path": str(tmp_path / "nope.pdf")})
    except ToolError:
        return
    raise AssertionError("Expected ToolError")


def test_read_directory(tmp_path):
    try:
        ReadDocumentTool().execute({"path": str(tmp_path)})
    except ToolError:
        return
    raise AssertionError("Expected ToolError")


def test_read_unsupported_extension(tmp_path):
    target = tmp_path / "photo.png"
    target.write_bytes(b"fake image")
    try:
        ReadDocumentTool().execute({"path": str(target)})
    except ToolError as exc:
        assert "Unsupported document type" in str(exc)
        return
    raise AssertionError("Expected ToolError")


def test_read_no_extension(tmp_path):
    target = tmp_path / "noext"
    target.write_bytes(b"data")
    try:
        ReadDocumentTool().execute({"path": str(target)})
    except ToolError as exc:
        assert "No file extension" in str(exc)
        return
    raise AssertionError("Expected ToolError")


def test_read_blank_path():
    try:
        ReadDocumentTool().execute({})
    except ToolError:
        return
    raise AssertionError("Expected ToolError")


# -- registry integration -------------------------------------------------

def test_registry_has_document_tool():
    registry = build_default_registry()
    assert registry.get("read_document") is not None