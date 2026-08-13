"""
Document intelligence tools (Phase 12).

`read_document` extracts readable text from common document formats so
JARVIS can read and answer questions about PDF, Word, Excel, PowerPoint,
CSV and plain-text files.

Supported extensions:
    * .txt .md .log .csv .json  - plain text
    * .pdf                      - pypdf (text layer; no OCR yet)
    * .docx                     - python-docx
    * .xlsx                     - openpyxl (values only, no formulas)
    * .pptx                     - python-pptx

Safety:
    * Paths are validated to exist and be files first.
    * Extraction is size-limited: results are truncated to a character
      budget so a huge document cannot flood the chat.
    * If the library for a format is not installed, the tool says exactly
      which one to `pip install` instead of crashing.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from tools.base import Tool, ToolError

_DEFAULT_LIMIT_CHARS = 6000
_MAX_LIMIT_CHARS = 100_000

_PLAIN_TEXT_EXTS = {".txt", ".md", ".log", ".csv", ".json", ".ini", ".yaml", ".yml"}
_SUPPORTED = "".join(
    sorted(_PLAIN_TEXT_EXTS | {".pdf", ".docx", ".xlsx", ".pptx"})
)


def _extract(text: str, limit: int) -> str:
    if not text or not text.strip():
        return ""
    snippet = text.strip()[:limit]
    if len(text) > limit:
        snippet += f"\n... [truncated: showing {limit} of {len(text)} chars]"
    return snippet


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ToolError(
            "Reading PDFs needs the 'pypdf' library. Run: pip install pypdf"
        )
    parts = []
    reader = PdfReader(str(path), strict=False)
    for index, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:  # degraded page must not kill the whole read
            text = ""
        if text.strip():
            parts.append(f"[page {index}]\n{text.strip()}")
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise ToolError(
            "Reading .docx files needs the 'python-docx' library. Run: pip install python-docx"
        )
    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ToolError(
            "Reading .xlsx files needs the 'openpyxl' library. Run: pip install openpyxl"
        )
    workbook = load_workbook(path, read_only=True, data_only=True)
    parts = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v).strip() for v in row]
                if any(cells):
                    parts.append(" | ".join(cells))
    finally:
        workbook.close()
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ToolError(
            "Reading .pptx files needs the 'python-pptx' library. Run: pip install python-pptx"
        )
    presentation = Presentation(str(path))
    parts = []
    for index, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in row.cells))
        if texts:
            parts.append(f"[slide {index}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


class ReadDocumentTool(Tool):
    name = "read_document"
    description = (
        "Extracts readable text from a document so you can answer questions "
        "about it. Supports " + _SUPPORTED + ". For plain text files, also "
        "accepts CSV and JSON as text."
    )
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to the document."},
            "max_chars": {
                "type": "integer",
                "description": "Optional limit on how much text to return.",
            },
        },
        "required": ["path"],
    }

    def execute(self, args: dict[str, Any]) -> str:
        path_raw = str(self._arg(args, "path", "") or "")
        max_chars = int(self._arg(args, "max_chars", 0) or 0)
        if max_chars <= 0:
            max_chars = _DEFAULT_LIMIT_CHARS
        max_chars = min(max_chars, _MAX_LIMIT_CHARS)

        path = Path(path_raw)
        if not path.exists():
            raise ToolError(f"File not found: {path}")
        if path.is_dir():
            raise ToolError(f"{path} is a directory, not a document.")
        ext = path.suffix.lower()
        if not ext:
            raise ToolError(f"No file extension to infer the format of: {path}")

        if ext in _PLAIN_TEXT_EXTS:
            try:
                text = path.read_text(encoding="utf-8-sig", errors="replace")
            except (OSError, PermissionError) as exc:
                raise ToolError(f"Could not read {path}: {exc}") from exc
            return _extract(text, max_chars) or f"{path} contains no text."

        extractor = {
            ".pdf": _extract_pdf,
            ".docx": _extract_docx,
            ".xlsx": _extract_xlsx,
            ".pptx": _extract_pptx,
        }.get(ext)
        if extractor is None:
            raise ToolError(f"Unsupported document type {ext!r}. Supported: {_SUPPORTED}.")

        try:
            text = extractor(path)
        except (OSError, PermissionError) as exc:
            raise ToolError(f"Could not read {path}: {exc}") from exc
        if not text.strip():
            return (
                f"{path} contains no extractable text "
                "(scanned/image PDFs need OCR, which is not supported yet)."
            )
        return _extract(text, max_chars)


def register_document_tools(registry) -> None:
    """Register the Phase 12 document-intelligence tools on a registry."""
    registry.register(ReadDocumentTool())